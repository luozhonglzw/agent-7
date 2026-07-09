"""PowerPoint parser using python-pptx.

Extracts text from slides and notes.

Usage::

    from app.core.parsers.pptx_parser import PPTXParser

    parser = PPTXParser()
    doc = parser.parse(Path("slides.pptx"))
"""

import logging
from pathlib import Path

from app.core.parsers.base import BaseParser, ParsedDocument, ParsedPage, ParserError

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """PPTX parser backed by python-pptx.

    Each slide becomes a :class:`ParsedPage`.  Speaker notes are
    appended to the slide content.
    """

    def supported_extensions(self) -> list[str]:
        """Return supported extensions.

        Returns:
            List containing ``".pptx"`` and ``".ppt"``.
        """
        return [".pptx", ".ppt"]

    def parse(self, file_path: Path) -> ParsedDocument:
        """Parse a PowerPoint file.

        Args:
            file_path: Absolute path to the PPTX file.

        Returns:
            Parsed document.

        Raises:
            ParserError: If python-pptx cannot open the file.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserError(
                detail="python-pptx is not installed",
                file_path=str(file_path),
            ) from exc

        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            raise ParserError(
                detail=f"Cannot open PPTX: {exc}",
                file_path=str(file_path),
            ) from exc

        pages: list[ParsedPage] = []
        full_text_parts: list[str] = []
        title = ""

        for idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []

            # Extract slide title if present
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = slide.shapes.title.text.strip()
                if slide_title:
                    parts.append(slide_title)
                    if not title:
                        title = slide_title

            # Extract text from all shapes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape is slide.shapes.title:
                    continue  # already handled
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

            # Speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")

            slide_text = "\n".join(parts)
            if slide_text.strip():
                pages.append(
                    ParsedPage(
                        page_number=idx,
                        content=slide_text,
                        heading=parts[0] if parts else None,
                        heading_level=1,
                    )
                )
                full_text_parts.append(slide_text)

        if not title:
            title = file_path.stem

        full_text = "\n\n".join(full_text_parts)

        logger.debug(
            "Parsed PPTX %s: %d slides, %d chars",
            file_path.name, len(pages), len(full_text),
        )

        return ParsedDocument(
            title=title,
            pages=pages,
            full_text=full_text,
            metadata={"slide_count": len(prs.slides)},
        )
